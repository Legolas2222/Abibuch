import pandas as pd
from pptx import Presentation

# Load Excel data
df = pd.read_excel("antworten.xlsx", sheet_name="data")

# Load the template
prs = Presentation("buchtemplate_single.pptx")
template_slide = prs.slides[0]

def replace_placeholders_in_shape(shape, data_row):
    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        # Join all runs' text
        full_text = ''.join(run.text for run in paragraph.runs)

        # Replace placeholders like {Name}, {Age}, etc.
        for key, value in data_row.items():
            print(f"Replacing {key} with {value}")
            placeholder = f"{{{key}}}"
            full_text = full_text.replace(placeholder, str(value))

        # Save style from first run
        if paragraph.runs:
            style_run = paragraph.runs[0]
            font = style_run.font
        else:
            font = None

        # Clear paragraph
        p = paragraph._element
        for r in list(p):
            p.remove(r)

        # Add one new run with updated text and preserved formatting
        new_run = paragraph.add_run()
        new_run.text = full_text

        # Apply style from original run
        if font:
            new_run.font.name = font.name
            new_run.font.size = font.size or Pt(18)  # Default size if none
            new_run.font.bold = font.bold
            new_run.font.italic = font.italic
            new_run.font.color.rgb = font.color.rgb if font.color and font.color.rgb else None

def create_slide_with_data(data_row):
    new_slide = prs.slides.add_slide(template_slide.slide_layout)
    for i, shape in enumerate(template_slide.shapes):
        if shape.has_text_frame:
            new_shape = new_slide.shapes[i]
            replace_placeholders_in_shape(new_shape, data_row)

# Generate one slide per row
for _, row in df.iterrows():
    create_slide_with_data(row)


# Save result
prs.save("formatted_presentation.pptx")
print("Presentation saved successfully with formatting preserved!")
