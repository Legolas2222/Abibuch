import pandas as pd
from pptx import Presentation

# Load Excel data
df = pd.read_excel("antworten.xlsx", sheet_name="data")

# Load the presentation
prs = Presentation("buchtemplate_single.pptx")
template_slide = prs.slides[0]

def update_textbox_text_preserving_formatting(shape, new_text):
    if not shape.has_text_frame:
        return

    # Use the first run and just update its text
    paragraph = shape.text_frame.paragraphs[0]

    # If there are no runs, create one
    if not paragraph.runs:
        run = paragraph.add_run()
    else:
        run = paragraph.runs[0]

    run.text = str(new_text)

    # Remove all other runs to prevent leftover text
    for old_run in paragraph.runs[1:]:
        paragraph._element.remove(old_run._element)

def create_slide_with_data(data_row):
    # Duplicate the slide layout
    new_slide = prs.slides.add_slide(template_slide.slide_layout)

    for i, shape in enumerate(template_slide.shapes):
        new_shape = new_slide.shapes[i]

        if shape.has_text_frame and shape.name in data_row:
            value = data_row[shape.name]
            update_textbox_text_preserving_formatting(new_shape, value)

# Apply for each row in Excel
for _, row in df.iterrows():
    create_slide_with_data(row)


# Save final result
prs.save("final_output.pptx")
print("Done! All formatting preserved, content inserted.")
